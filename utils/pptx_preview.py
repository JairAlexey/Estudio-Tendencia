import io
from pptx import Presentation
import streamlit as st

def extract_pptx_content(pptx_data):
    """Extrae el contenido del PPTX y lo convierte a HTML"""
    try:
        # Crear archivo temporal en memoria
        pptx_stream = io.BytesIO(pptx_data)
        prs = Presentation(pptx_stream)
        
        slides_html = []
        for i, slide in enumerate(prs.slides, 1):
            # Contenedor para cada diapositiva
            slide_content = [f'<div style="margin-bottom: 2rem; padding: 2rem; border: 1px solid #ddd; border-radius: 8px;">']
            slide_content.append(f'<h3 style="color: #333;">Diapositiva {i}</h3>')
            
            # Extraer texto de cada shape en la diapositiva
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Formatear el texto dependiendo del nivel
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                # Determinar el estilo basado en el nivel
                                if paragraph.level == 0:
                                    slide_content.append(f'<p style="font-size: 1.2em; font-weight: bold; color: #2c5282;">{text}</p>')
                                else:
                                    slide_content.append(f'<p style="margin-left: {paragraph.level * 20}px; color: #444;">{text}</p>')
            
            slide_content.append('</div>')
            slides_html.append('\n'.join(slide_content))
        
        return '\n'.join(slides_html)

    except Exception as e:
        st.error(f"Error al procesar PPTX: {str(e)}")
        return None

def mostrar_preview_pptx(file_data):
    """Muestra una previsualización del PPTX como HTML estilizado"""
    try:
        if isinstance(file_data, memoryview):
            file_data = file_data.tobytes()
        
        # Estilo general para la previsualización
        st.markdown("""
            <style>
            .pptx-preview {
                background: white;
                padding: 2rem;
                border-radius: 10px;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            }
            </style>
        """, unsafe_allow_html=True)
        
        html_content = extract_pptx_content(file_data)
        if html_content:
            st.markdown(
                f'<div class="pptx-preview">{html_content}</div>',
                unsafe_allow_html=True
            )
        else:
            st.error("No se pudo procesar la presentación")
            
    except Exception as e:
        st.error(f"Error al mostrar la presentación: {str(e)}")
        st.error(f"Error al mostrar la presentación: {str(e)}")
        st.download_button(
            "Descargar PPTX",
            data=file_data,
            file_name="presentacion.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
