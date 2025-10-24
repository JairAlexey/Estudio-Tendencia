import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from conexion import conn, cursor
from pptx import Presentation
from tools.generar_mapa_latam import generar_mapa_latam

# Obtener el directorio raíz del proyecto
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


# ==========================================================
#   FUNCIONES AUXILIARES
# ==========================================================
def obtener_datos_solicitud_por_proyecto(proyecto_id):
    cursor.execute('''
        SELECT nombre_proponente, duracion, modalidad, nombre_programa,
               facultad_propuesta, facultad_proponente, cargo_proponente
        FROM datos_solicitud
        WHERE proyecto_id = %s
    ''', (proyecto_id,))
    row = cursor.fetchone()
    if row:
        return {
            "nombre_proponente": row[0],
            "duracion": row[1],
            "modalidad": row[2],
            "nombre_programa": row[3],
            "facultad_propuesta": row[4],
            "facultad_proponente": row[5],
            "cargo_proponente": row[6]
        }
    else:
        return None


def reemplazar_texto_preservando_formato(shape, nuevo_valor):
    """
    Inserta o reemplaza el texto después de ':' sin borrar formato previo.
    Si no hay texto después de los ':', simplemente agrega el nuevo valor.
    Maneja colores y formatos ausentes sin errores.
    """
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame
    texto_actual = shape.text.strip()

    if ":" not in texto_actual:
        return

    etiqueta, _, valor_existente = texto_actual.partition(":")
    etiqueta = etiqueta.strip()
    valor_existente = valor_existente.strip()
    nuevo_texto = f" {nuevo_valor}"

    reemplazado = False

    for paragraph in text_frame.paragraphs:
        if not paragraph.runs:
            continue

        for run in paragraph.runs:
            if valor_existente and valor_existente in run.text:
                run.text = run.text.replace(valor_existente, nuevo_texto)
                reemplazado = True
                break
        if reemplazado:
            break

    # Si no se reemplazó (porque no había valor existente)
    if not reemplazado:
        p = text_frame.paragraphs[-1]
        run_base = p.runs[-1] if p.runs else p.add_run()
        nuevo_run = p.add_run()

        # Copiar formato de forma segura
        nuevo_run.font.bold = getattr(run_base.font, "bold", None)
        nuevo_run.font.size = getattr(run_base.font, "size", None)
        nuevo_run.font.name = getattr(run_base.font, "name", None)
        try:
            if run_base.font.color and hasattr(run_base.font.color, "rgb"):
                nuevo_run.font.color.rgb = run_base.font.color.rgb
        except Exception:
            pass

        nuevo_run.text = nuevo_texto


# ==========================================================
#   GENERAR REPORTE DE VIABILIDAD
# ==========================================================
def generar_reporte(proyecto_id, viabilidad=None):
    datos = obtener_datos_solicitud_por_proyecto(proyecto_id)
    if not datos:
        print("No se encontraron datos para la solicitud.")
        return

    # Seleccionar plantilla según viabilidad
    if viabilidad is not None:
        if viabilidad <= 60:
            template_path = os.path.join(PROJECT_ROOT, 'files/base/NoViable.pptx')
        elif viabilidad <= 70:
            template_path = os.path.join(PROJECT_ROOT, 'files/base/Revision.pptx')
        else:
            template_path = os.path.join(PROJECT_ROOT, 'files/base/Viable.pptx')
    else:
        template_path = os.path.join(PROJECT_ROOT, 'files/base/Viable.pptx')

    prs = Presentation(template_path)
    slide = prs.slides[1]

    shape_map = {
        1: "nombre_proponente",
        3: "duracion",
        4: "modalidad",
        5: "nombre_programa",
        10: "facultad_propuesta",
        11: "facultad_proponente",
        14: "cargo_proponente"
    }

    for idx, col in shape_map.items():
        try:
            shape = slide.shapes[idx]
            reemplazar_texto_preservando_formato(shape, datos[col])
        except Exception as e:
            print(f"Error procesando shape {idx}: {e}")

    # Reemplazar viabilidad en el slide 3
    if viabilidad is not None:
        try:
            slide_img = prs.slides[3]
            shape_viabilidad = slide_img.shapes[11]
            reemplazar_texto_preservando_formato(shape_viabilidad, f"{viabilidad:.1f}%")
        except Exception as e:
            print(f"No se pudo reemplazar la viabilidad en el slide 3: {e}")

    # Reemplazar imagen radar (slide 3)
    try:
        slide_img = prs.slides[3]
        for shape in slide_img.shapes:
            if shape.shape_type == 13:  # Picture
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                slide_img.shapes._spTree.remove(shape._element)
                ruta_img = os.path.join(PROJECT_ROOT, f"files/imagenes/grafico_radar_{proyecto_id}.png")
                slide_img.shapes.add_picture(ruta_img, left, top, width, height)
                break
    except Exception as e:
        print(f"No se pudo reemplazar la imagen en el slide 3: {e}")

    nombre_archivo = f"{datos['nombre_programa'].replace(' ', '_')}_viabilidad.pptx"
    output_path = os.path.join(PROJECT_ROOT, 'files/presentaciones', nombre_archivo)
    prs.save(output_path)
    print(f"Reporte guardado en: {output_path}")
    return output_path


# ==========================================================
#   GENERAR REPORTE DE MERCADO
# ==========================================================
def generar_reporte_mercado(proyecto_id):
    """
    Genera un reporte de investigación de mercado respetando los formatos
    y reemplazando solo el contenido después de los ':'.
    Además, genera el gráfico de mapa Latam y lo inserta en el shape 6.
    """
    datos = obtener_datos_solicitud_por_proyecto(proyecto_id)
    if not datos:
        print("No se encontraron datos para la solicitud.")
        return

    template_path = os.path.join(PROJECT_ROOT, 'files/base/InvestigacionMercados.pptx')
    prs = Presentation(template_path)

    slide_index = 1
    if slide_index >= len(prs.slides):
        slide_index = len(prs.slides) - 1
    slide = prs.slides[slide_index]

    # Según tu nuevo inspect, estos son los índices correctos
    shape_map = {
        1: "nombre_proponente",
        3: "duracion",
        4: "modalidad",
        5: "nombre_programa",
        10: "facultad_propuesta",
        11: "facultad_proponente",
        14: "cargo_proponente"
    }

    for idx, col in shape_map.items():
        try:
            shape = slide.shapes[idx]
            reemplazar_texto_preservando_formato(shape, datos[col])
        except Exception as e:
            print(f"Error procesando shape {idx}: {e}")

 
    # --- Insertar aptitudes en la tabla del slide 7 ---
    try:
        slide_aptitudes = prs.slides[7]
        # Buscar el shape tipo tabla (shape_type == 19)
        tabla_shape = None
        for shape in slide_aptitudes.shapes:
            if shape.shape_type == 19:
                tabla_shape = shape
                break
        if tabla_shape:
            table = tabla_shape.table
            # Extraer aptitudes de la base de datos (traer también cantidad)
            cursor.execute("""
                SELECT nombre, porcentaje, cantidad, ubicacion
                FROM linkedin_aptitudes
                WHERE proyecto_id = %s AND ubicacion IN ('Ecuador', 'América Latina')
            """, (proyecto_id,))
            rows = cursor.fetchall()
            # Limpiar el símbolo % para ordenamiento
            def limpiar_porcentaje(val):
                try:
                    return float(str(val).replace('%', '').replace(',', '.').strip())
                except Exception:
                    return 0.0
            def limpiar_cantidad(val):
                try:
                    return float(str(val).replace(',', '').strip())
                except Exception:
                    return 0.0
            # Separar por ubicación y ordenar por porcentaje (desc), luego cantidad (desc)
            aptitudes_ecuador = sorted(
                [r for r in rows if r[3] == "Ecuador"],
                key=lambda x: (limpiar_porcentaje(x[1]), limpiar_cantidad(x[2])),
                reverse=True
            )[:10]
            aptitudes_latam = sorted(
                [r for r in rows if r[3] == "América Latina"],
                key=lambda x: (limpiar_porcentaje(x[1]), limpiar_cantidad(x[2])),
                reverse=True
            )[:10]
            # Rellenar la tabla (asume estructura: fila 2 a 11 para datos)
            for i in range(10):
                # Ecuador
                if i < len(aptitudes_ecuador):
                    nombre_ec = aptitudes_ecuador[i][0]
                    porcentaje_ec = str(aptitudes_ecuador[i][1])
                else:
                    nombre_ec = ""
                    porcentaje_ec = ""
                # Latam
                if i < len(aptitudes_latam):
                    nombre_lat = aptitudes_latam[i][0]
                    porcentaje_lat = str(aptitudes_latam[i][1])
                else:
                    nombre_lat = ""
                    porcentaje_lat = ""
                # Fila de la tabla (fila 2+i)
                row_idx = 2 + i
                # Ecuador nombre
                cell_ec = table.cell(row_idx, 0)
                for paragraph in cell_ec.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = nombre_ec
                        break
                # Ecuador porcentaje
                cell_pct_ec = table.cell(row_idx, 1)
                for paragraph in cell_pct_ec.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = porcentaje_ec
                # Latam nombre
                cell_lat = table.cell(row_idx, 3)
                for paragraph in cell_lat.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = nombre_lat
                        break
                # Latam porcentaje
                cell_pct_lat = table.cell(row_idx, 4)
                for paragraph in cell_pct_lat.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = porcentaje_lat
            print("Aptitudes insertadas en la tabla del slide 7.")
        else:
            print("No se encontró tabla en el slide 7 para insertar aptitudes.")
    except Exception as e:
        print(f"Error insertando aptitudes en la tabla del slide 7: {e}")

    # --- Actualizar tabla de evaluación en slide 9 ---
    try:
        actualizar_tabla_evaluacion_slide9(proyecto_id, prs)
    except Exception as e:
        print(f"Error actualizando tabla de evaluación en slide 9: {e}")

    # --- Reemplazar imagen del mapa en slide 5, shape 6 ---
    try:
        # Generar el mapa de Latam
        ruta_mapa = os.path.join(PROJECT_ROOT, f"files/imagenes/mapa_latam_{proyecto_id}.png")
        generar_mapa_latam(proyecto_id, ruta_mapa)
        
        # Acceder al slide 5 y reemplazar la imagen en shape 6
        slide_mapa = prs.slides[5]
        shape_imagen = slide_mapa.shapes[6]
        
        if shape_imagen.shape_type == 13:  # Picture
            # Guardar las propiedades de la imagen original
            left, top, width, height = shape_imagen.left, shape_imagen.top, shape_imagen.width, shape_imagen.height
            
            # Eliminar la imagen original
            slide_mapa.shapes._spTree.remove(shape_imagen._element)
            
            # Agregar la nueva imagen con las mismas propiedades
            slide_mapa.shapes.add_picture(ruta_mapa, left, top, width, height)
            print(f"Imagen del mapa reemplazada en slide 5, shape 6: {ruta_mapa}")
        else:
            print("El shape 6 en slide 5 no es una imagen")
    except Exception as e:
        print(f"Error reemplazando imagen del mapa en slide 5: {e}")

    nombre_archivo = f"{datos['nombre_programa'].replace(' ', '_')}_mercado.pptx"
    output_path = os.path.join(PROJECT_ROOT, 'files/presentaciones', nombre_archivo)
    prs.save(output_path)
    print(f"Reporte de mercado guardado en: {output_path}")
    return output_path


def actualizar_tabla_evaluacion_slide9(proyecto_id, prs):
    """
    Actualiza la tabla de evaluación en el slide 9, shape 0 (tabla),
    modificando solo las columnas de Presencialidad y Virtualidad con los valores
    obtenidos de grafico_radar_datos para el proyecto_id.
    """
    # Obtener los valores desde la base de datos
    cursor.execute("""
        SELECT valor_busqueda, valor_competencia_presencialidad, valor_competencia_virtualidad, valor_linkedin, valor_mercado
        FROM grafico_radar_datos
        WHERE proyecto_id = %s
    """, (proyecto_id,))
    row = cursor.fetchone()
    if not row:
        print(f"No se encontraron datos en grafico_radar_datos para proyecto_id={proyecto_id}")
        return

    valor_busqueda = float(row[0]) if row[0] is not None else 0.0
    valor_competencia_presencialidad = float(row[1]) if row[1] is not None else 0.0
    valor_competencia_virtualidad = float(row[2]) if row[2] is not None else 0.0
    valor_linkedin = float(row[3]) if row[3] is not None else 0.0
    valor_mercado = float(row[4]) if row[4] is not None else 0.0

    # Calcular totales
    total_presencialidad = round(valor_busqueda + valor_competencia_presencialidad + valor_linkedin + valor_mercado, 2)
    total_virtualidad = round(valor_busqueda + valor_competencia_virtualidad + valor_linkedin + valor_mercado, 2)

    # Acceder al slide 9 y shape 0 (tabla)
    slide_eval = prs.slides[9]
    tabla_shape = slide_eval.shapes[0]
    table = tabla_shape.table

    # Fila 3: Búsqueda Web
    cell_presencial = table.cell(3, 2)
    cell_virtual = table.cell(3, 3)
    for paragraph in cell_presencial.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = f"{valor_busqueda:.0f}%"
    for paragraph in cell_virtual.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = f"{valor_busqueda:.0f}%"

    # Fila 4: LinkedIN
    cell_presencial = table.cell(4, 2)
    cell_virtual = table.cell(4, 3)
    for paragraph in cell_presencial.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = f"{valor_linkedin:.0f}%"
    for paragraph in cell_virtual.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = f"{valor_linkedin:.0f}%"

    # Fila 5: Competencia
    cell_presencial = table.cell(5, 2)
    cell_virtual = table.cell(5, 3)
    for paragraph in cell_presencial.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = f"{valor_competencia_presencialidad:.0f}%"
    for paragraph in cell_virtual.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = f"{valor_competencia_virtualidad:.0f}%"

    # Fila 6: Actividades Económicas
    cell_presencial = table.cell(6, 2)
    cell_virtual = table.cell(6, 3)
    for paragraph in cell_presencial.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = f"{valor_mercado:.0f}%"
    for paragraph in cell_virtual.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = f"{valor_mercado:.0f}%"

    # Fila 7: Total (suma de todas las anteriores)
    cell_presencial = table.cell(7, 2)
    cell_virtual = table.cell(7, 3)
    for paragraph in cell_presencial.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = f"{total_presencialidad:.0f}%"
    for paragraph in cell_virtual.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = f"{total_virtualidad:.0f}%"

    print("Tabla de evaluación en slide 9 actualizada correctamente.")


# ==========================================================
#   MAIN (prueba rápida)
# ==========================================================
if __name__ == "__main__":
    proyecto_id = 1  # Cambia esto según tu caso
    generar_reporte(proyecto_id)
    generar_reporte_mercado(proyecto_id)