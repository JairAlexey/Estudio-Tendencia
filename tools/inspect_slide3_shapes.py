from pptx import Presentation

# Ruta de la plantilla
TEMPLATE_PATH = 'db/Plantilla para automatización de tendencias.pptx'

# Cargar la presentación
prs = Presentation(TEMPLATE_PATH)

# Acceder a la diapositiva 3 (índice 2)
slide = prs.slides[2]

print('Textos de cada shape en la diapositiva 3:')
for i, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        print(f"Shape {i}: {shape.text}")
