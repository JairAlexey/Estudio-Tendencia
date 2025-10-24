from pptx import Presentation

# Ruta de la plantilla
TEMPLATE_PATH = '../files/base/InvestigacionMercados.pptx'

# Cargar la presentación
prs = Presentation(TEMPLATE_PATH)


slide = prs.slides[5]

print('Textos de cada shape:')
for i, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        print(f"Shape {i}: {shape.text}")
    elif shape.shape_type == 13:  # 13 es MSO_SHAPE_TYPE.PICTURE
        print(f"Shape {i}: [Imagen] - Tamaño: {shape.width}x{shape.height}, Posición: ({shape.left},{shape.top})")
    elif shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
        print(f"Shape {i}: [Tabla]")
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            row_text = []
            for cell in row.cells:
                row_text.append(cell.text)
            print(f"  Fila {row_idx}: {row_text}")
